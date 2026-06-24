#!/usr/bin/env python3
"""
Create PowerPoint presentation for Manager Discussion: IDP & Career/Compensation Alignment
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(68, 114, 196)
ACCENT_ORANGE = RGBColor(237, 125, 49)
DARK_GRAY = RGBColor(68, 68, 68)
TEXT_COLOR = RGBColor(51, 51, 51)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_ORANGE

def add_content_slide(prs, title, bullets, title_color=DARK_BLUE):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title bar
    title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = title_color
    title_box.line.color.rgb = title_color
    
    # Title text
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_bottom = Inches(0.1)
    title_frame.margin_top = Inches(0.1)
    title_frame.margin_left = Inches(0.3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = DARK_BLUE
    title_box.line.color.rgb = DARK_BLUE
    
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Add table
    rows_count = len(rows) + 1
    cols_count = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.5)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Set column widths
    for col_idx, header in enumerate(headers):
        table_shape.columns[col_idx].width = Inches(width.inches / cols_count)
    
    # Header row
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table_shape.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
            
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT_COLOR

# SLIDE 1: Title Slide
add_title_slide(prs, "IDP Review & Career Discussion", "Title & Compensation Alignment")

# SLIDE 2: Agenda
add_content_slide(prs, "Today's Discussion", [
    "• IDP Overview (Goals 1-7)",
    "• Development Plan through Jun 2027",
    "• Resources & Support Needed",
    "• Title & Compensation Alignment (PRIMARY)",
    "• Next Steps & Timeline"
])

# SLIDE 3: Who I Am
add_content_slide(prs, "Your Unique Value: Full-Stack Systems Engineer", [
    "• Only engineer on team with expertise across entire system lifecycle",
    "• Hardware Design → Firmware → Test Automation → Manufacturing → Certification → Cloud → Field Support",
    "• Most peers specialize in ONE domain; you own ALL of them",
    "• Rare skill set: end-to-end ownership, autonomous execution, cross-functional leadership",
    "• KEY MESSAGE: You solve problems no one else on the team can"
])

# SLIDE 4: Technical Expertise Summary
headers = ["Domain", "Level", "Key Examples"]
rows = [
    ["Hardware Design (Altium)", "5 Expert", "M1, Mercury boards, production fixtures"],
    ["Firmware Development", "5 Expert", "Bare metal, real-time, hardware abstraction"],
    ["Test Fixture Systems", "5 Expert", "M1 & Mercury systems, 0.1% accuracy"],
    ["Test Automation & Data", "4 Advanced", "2.5M swipes stress platform, databases"],
    ["Cloud & Microservices", "4 Advanced", "Elements gateway, Azure, OTA updates"],
    ["Manufacturing", "5 Expert", "Design-for-manufacturability, yield optimization"],
    ["Certification & Compliance", "4 Advanced", "UL/TUV/Intertek on-site participation"],
    ["Cross-Functional Leadership", "5 Expert", "Field support, customer escalations, mentoring"]
]
add_table_slide(prs, "8 Technical Competencies", headers, rows)

# SLIDE 5: Real Work Examples
add_content_slide(prs, "Evidence of Full-Stack Expertise", [
    "Elements Gateway Design:",
    "  - Designed entire system from ground up (hardware + firmware + cloud)",
    "  - Bridges legacy hardware with modern Azure platform",
    "",
    "Test Fixture Systems:",
    "  - Complete design: M1 test board, Mercury test board, production fixtures",
    "  - Hardware (PCB), firmware (real-time control), software (automation)",
    "",
    "Manufacturing & Field Support:",
    "  - On-site factory collaboration, design-for-manufacturability, yield optimization",
    "  - Certification labs (UL, TUV, Intertek), customer escalations, field diagnostics"
])

# SLIDE 6: Title Issue
add_content_slide(prs, "Current Title Misrepresents Actual Expertise", [
    "What the Title Says:",
    "  - \"Software Leader\" = pure software specialty OR management role",
    "",
    "What Reality Shows:",
    "  - Principal-level Systems Engineer",
    "  - Hardware + Firmware + Software + Cloud + Manufacturing + Compliance",
    "  - YOU ARE THE ONLY ENGINEER ON THE TEAM WITH THIS BREADTH",
    "",
    "Why It Matters:",
    "  - Internal equity: misrepresents scope to HR",
    "  - External mobility: doesn't reflect actual specialty",
    "  - Compensation: \"Software Leader\" rates are lower than \"Principal Engineer, Systems Architecture\""
])

# SLIDE 7: Market Analysis
add_content_slide(prs, "Compensation Market Reality", [
    "Current Situation:",
    "  - Base Salary: $185,000 (as \"Software Leader\")",
    "  - Experience: 34 years total, 28+ years post-immigration",
    "  - Career Stage: Principal-level expert with full-stack mastery",
    "",
    "Market Data (Principal Engineer, Systems/Hardware Focus):",
    "  - 15% adjustment: $212,750 (Conservative)",
    "  - 20% adjustment: $222,000 (Market-Aligned)",
    "  - 25% adjustment: $231,250 (Competitive)",
    "  - Industry standard: $240,000 - $280,000+",
    "",
    "KEY: You're likely 15-25% below market for your actual role"
])

# SLIDE 8: Business Case
add_content_slide(prs, "Why This Adjustment Makes Business Sense", [
    "1. Rare Skill Set: Only engineer on team with complete full-stack expertise",
    "",
    "2. Scope of Responsibility: End-to-end ownership across hardware→cloud→field",
    "",
    "3. Travel Commitment: Customer sites, factories, certification labs, weekend availability",
    "",
    "4. Strategic Value: Solves problems no one else can, bridges hardware and cloud domains",
    "",
    "5. Market Reality: Principal Engineer roles command 15-25% premium over \"Software Leader\""
])

# SLIDE 9: IDP Goals Overview
headers = ["Goal", "Target", "By Date"]
rows = [
    ["1. Manufacturing Resilience", "Advanced→Expert", "Dec 31, 2026"],
    ["2. Certification Leadership", "Advanced→Expert", "Jun 30, 2027"],
    ["3. Data-Driven Design", "Proficient→Advanced", "Dec 31, 2026"],
    ["4. Design Standards", "Advanced→Expert", "Dec 31, 2026"],
    ["5. Cloud Architecture", "Advanced→Expert", "Dec 31, 2026"],
    ["6. Cross-Product Platform", "Advanced→Expert", "Jun 30, 2027"],
    ["7. Mentoring & Leadership", "Advanced→Expert", "Ongoing"]
]
add_table_slide(prs, "7 Development Goals (May 2026 - Jun 2027)", headers, rows)

# SLIDE 10: Checkpoints & Accountability
add_content_slide(prs, "Quarterly Checkpoints & Measurable Progress", [
    "Checkpoint Dates:",
    "  - Aug 31, 2026: Cycle 1 progress review",
    "  - Dec 31, 2026: Cycle 2 progress review",
    "  - Jun 30, 2027: Full cycle evaluation & progression to Expert level",
    "",
    "Success Measures:",
    "  ✓ Business Impact: Reduced field escalation time, faster diagnostics",
    "  ✓ Team Capability: Engineers solving cross-domain problems independently",
    "  ✓ Documentation: Standards adopted across teams",
    "  ✓ Individual Growth: Progression from Advanced (4) toward Expert (5)"
])

# SLIDE 11: Resources Needed
add_content_slide(prs, "Support Required for IDP Goals 1-7", [
    "Time Allocation:",
    "  - 60-70% hands-on engineering (design, firmware, cloud work)",
    "  - 20-30% mentoring, standards documentation, architecture",
    "  - 10-20% cross-functional alignment",
    "",
    "Tools & Infrastructure:",
    "  - Cloud platform tools (deployment infrastructure)",
    "  - Hardware diagnostic tools (~$5K for specialized equipment)",
    "  - Documentation system (Confluence, GitHub)",
    "",
    "Training & Development:",
    "  - Cloud architecture certification (AWS/Azure)",
    "  - Advanced DevOps practices, Kubernetes/containerization"
])

# SLIDE 12: Recommended Title
add_content_slide(prs, "Proposed Title: \"Principal Engineer, Systems Architecture\"", [
    "Why This Title?",
    "",
    "Accuracy:",
    "  - Reflects actual discipline (Systems Engineer, not Software)",
    "  - Describes scope (Architecture across hardware-cloud)",
    "  - Aligns with Honeywell IC technical track",
    "",
    "Market Alignment:",
    "  - Accurate for industry mobility",
    "  - Reflects Principal-level expertise",
    "  - Supports compensation adjustments"
])

# SLIDE 13: Compensation Recommendation
add_content_slide(prs, "Proposed Compensation Adjustment", [
    "Recommended Range:",
    "  - Target: $220,000 - $231,000 (20-25% adjustment)",
    "  - Rationale: Principal-level IC with rare full-stack systems engineer expertise",
    "  - Supports: Title change to \"Principal Engineer, Systems Architecture\"",
    "",
    "Timeline:",
    "  - Discuss alignment today",
    "  - HR review & approval (typically 1-2 weeks)",
    "  - Implementation in next pay cycle",
    "",
    "Next Steps from Manager:",
    "  1. Validate market analysis with HR",
    "  2. Present business case to compensation/HR team",
    "  3. Discuss implementation timeline"
])

# SLIDE 14: Summary & Ask
add_content_slide(prs, "What I'm Asking For", [
    "1. TITLE CHANGE:",
    "   From: \"Software Leader\" → To: \"Principal Engineer, Systems Architecture\"",
    "",
    "2. COMPENSATION ADJUSTMENT:",
    "   Range: $220,000 - $231,000 (20-25% adjustment)",
    "   Only engineer on team with this breadth",
    "",
    "3. SUPPORT FOR IDP GOALS 1-7:",
    "   - 60-70% hands-on engineering time",
    "   - Tools, training, budget as outlined",
    "   - Quarterly checkpoints (Aug 31, Dec 31, Jun 30)"
])

# SLIDE 15: Business Impact
add_content_slide(prs, "What This Enables for Honeywell", [
    "Retention:",
    "  - Recognizes and properly compensates rare expertise",
    "  - Supports career progression at Principal level",
    "",
    "Knowledge & Standards:",
    "  - Principal-level systems engineering standards across portfolio",
    "  - Cross-product platform standardization (30% time-to-market improvement)",
    "",
    "Competitive Advantage:",
    "  - Gateway design bridges hardware and cloud (rare)",
    "  - Full-stack problem-solving, field expertise supporting customer success"
])

# SLIDE 16: Next Steps
add_content_slide(prs, "Path Forward", [
    "Immediate Actions (This Week):",
    "  - Agree on title change and compensation adjustment approach",
    "  - Identify HR/compensation contacts for discussion",
    "",
    "Next 1-2 Weeks:",
    "  - HR review and market validation",
    "  - Discuss implementation timeline",
    "",
    "Ongoing:",
    "  - Monthly check-ins on IDP progress (Goals 1-7)",
    "  - Quarterly checkpoints (Aug 31, Dec 31, Jun 30)"
])

# SLIDE 17: Questions & Discussion
add_content_slide(prs, "Let's Discuss", [
    "Key Discussion Points:",
    "",
    "On Market Data:",
    "  Q: How confident are you in the 15-25% range?",
    "  A: Industry benchmarks for Principal Engineer (systems focus) roles; Honeywell's own IC bands",
    "",
    "On Title: Will this require HR approval? → Yes, HR will review and approve",
    "",
    "On Compensation: What's the business justification? → Rarity of expertise, scope, market alignment",
    "",
    "On IDP Goals: How realistic is this timeline? → 12-18 months is standard for Advanced→Expert progression"
])

# SLIDE 18: Thank You
add_title_slide(prs, "Thank You", "Questions?")

# Save presentation
output_path = r"c:\Users\H573359\myGithub\M1Combined.worktrees\agents-design-review-cleanup\m1tfc\Manager_Discussion_Presentation.pptx"
prs.save(output_path)
print(f"✓ PowerPoint presentation created: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"\nPresentation ready for your manager discussion!")
